"""redesign_pptx.py
Redesign all presentation slides: image placeholders, minimal text, MPJPE fix.
Usage:  python -m scripts.redesign_pptx
"""
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

PPTX  = 'Documentation/Piano_Finger_Tracking_Presentation.pptx'
PLOTS = Path('data/plots')
ROOT  = Path('.')

prs = Presentation(PPTX)
W, H = prs.slide_width, prs.slide_height   # EMU
SL   = prs.slides


# ── helpers ──────────────────────────────────────────────────────────────────

def clear(slide):
    sp = slide.shapes._spTree
    for el in list(sp):
        if el.tag.rsplit('}', 1)[-1] in ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'):
            sp.remove(el)

def fix_mjmpe(slide):
    for s in slide.shapes:
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    r.text = r.text.replace('MJMPE', 'MPJPE')

def ttl(slide, text, size=30):
    tb = slide.shapes.add_textbox(In(0.45), In(0.18), W - In(0.9), In(0.72))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text, r.font.size, r.font.bold = text, Pt(size), True
    r.font.color.rgb = RGBColor(255, 255, 255)

def txt(slide, lines, l, t, w, h, size=13, bullet=True, color=(210, 215, 230)):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = ('• ' if bullet else '') + ln
        r.font.size, r.font.color.rgb = Pt(size), RGBColor(*color)

def box(slide, l, t, w, h, label, sub=''):
    """Dark navy placeholder rectangle."""
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(30, 42, 65)
    s.line.color.rgb = RGBColor(75, 100, 155); s.line.width = Pt(1.5)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size, r.font.bold = Pt(10), True
    r.font.color.rgb = RGBColor(140, 165, 215)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size, r2.font.italic = Pt(9), True
        r2.font.color.rgb = RGBColor(115, 140, 190)

def lbl(slide, text, l, t, w, size=10, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(l, t, w, In(0.32))
    p = tb.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size, r.font.italic = Pt(size), True
    r.font.color.rgb = RGBColor(170, 182, 210)

def pic(slide, path, l, t, w=None, h=None, fallback_label=None):
    fp = Path(path)
    if fp.exists():
        return slide.shapes.add_picture(str(fp), l, t, width=w, height=h)
    box(slide, l, t, w or In(5), h or In(4), fallback_label or fp.name)


# ── card grid helpers ─────────────────────────────────────────────────────────

def card_row(slide, items, y_img, img_h, y_lbl, margin_l=In(0.5), gap=In(0.25)):
    """
    items = list of (label, sublabel)  — draws image placeholders in a row.
    Returns list of (left_x, card_width) for each card.
    """
    n = len(items)
    avail = W - margin_l * 2
    cw = (avail - gap * (n - 1)) / n
    positions = []
    for i, (label, sub) in enumerate(items):
        lx = margin_l + i * (cw + gap)
        box(slide, lx, y_img, cw, img_h, label, sub)
        positions.append((lx, cw))
    return positions, cw


# ── slide 4 (index 3): Participant Study Overview ─────────────────────────────

s4 = SL[3]
clear(s4)
ttl(s4, 'Participant Study Overview')
pic(s4, PLOTS / '00_participant_composition.png',
    In(0.8), In(1.0), w=In(11.7),
    fallback_label='00_participant_composition.png')
txt(s4, [
    '38 participants  ·  types II–V  ·  dim & indoor lighting  ·  7.7–12.2 cm hand size',
    'Most participants: Fitzpatrick III–IV  ·  majority in dim lighting  ·  hand sizes > 10 cm',
], In(0.5), In(6.0), W - In(1.0), In(1.3), size=13)


# ── slide 6 (index 5): Literature Review ─────────────────────────────────────

s6 = SL[5]
clear(s6)
ttl(s6, 'Literature Review')

papers = [
    ('📷  Kinect depth sensor\nimage placeholder',  'Han et al. (2013)'),
    ('🤚  MediaPipe 21-landmark\nmodel diagram',    'Zhang et al. (2020)'),
    ('🕺  OpenPose PAF body\npose skeleton',        'Cao et al. (2019)'),
    ('🎹  AR piano overlay\non real keyboard',      'Hiranaka et al. (2022)'),
]
pos, cw = card_row(s6, papers, y_img=In(1.05), img_h=In(3.5), y_lbl=In(4.6))
for (lx, _), (_, cite) in zip(pos, papers):
    lbl(s6, cite, lx, In(4.62), cw, size=11)

txt(s6, [
    'Molchanov et al. (2016): single-camera CNN landmark detection — no depth hardware needed',
    'Liang et al. (2016): RGB-D virtual piano (>90% accuracy but requires depth camera)',
    'Lugaresi et al. (2019): temporal smoothing reduces jitter but adds latency on fast movements',
    'Pu et al. (2023): real-world occlusion & rapid motion reduce accuracy more than benchmarks',
], In(0.5), In(5.05), W - In(1.0), In(2.2), size=11.5)


# ── slide 8 (index 7): Hardware Components ───────────────────────────────────

s8 = SL[7]
clear(s8)
ttl(s8, 'Hardware Components')

hw = [
    ('Casio CTK-2400\n61-key piano',    'MIDI ground truth'),
    ('Google Pixel 7\nUSB Type-C',      'Overhead camera'),
    ('Velbon tripod\nadjustable',       'Fixed overhead mount'),
    ('MSI laptop\ni7-14th / RTX 5070',  'CUDA processing'),
    ('~90° overhead\nsetup',             'Overhead view of both hands'),
]

# Last item: use actual session_example.jpg if available
n = len(hw)
avail = W - In(0.7)
gap = In(0.2)
cw = (avail - gap * (n - 1)) / n
lx0 = In(0.35)

for i, (label, sub) in enumerate(hw):
    lx = lx0 + i * (cw + gap)
    if i == n - 1 and (ROOT / 'session_example.jpg').exists():
        pic(s8, ROOT / 'session_example.jpg', lx, In(1.05), w=cw, h=In(3.7))
    else:
        box(s8, lx, In(1.05), cw, In(3.7), label, sub)
    lbl(s8, label.split('\n')[0], lx, In(4.82), cw, size=10)
    lbl(s8, sub, lx, In(5.12), cw, size=9)


# ── slide 9 (index 8): Software Components ───────────────────────────────────

s9 = SL[8]
clear(s9)
ttl(s9, 'Software Components & Libraries')

sw_logos = [
    ('Python 3.10\nlogo', 'Compatibility with\nMediaPipe & CUDA'),
    ('OpenCV (cv2)\nlogo', 'Camera capture\nhomography & overlay'),
    ('MediaPipe Hands\narchitecture', '21 keypoints · real-time\nno GPU required'),
]
pos9, cw9 = card_row(s9, [(l, s) for l, s in sw_logos],
                     y_img=In(1.05), img_h=In(3.0), y_lbl=In(0), gap=In(0.3))
for (lx, _), (lbl_t, _) in zip(pos9, sw_logos):
    lbl(s9, lbl_t.split('\n')[0], lx, In(4.12), cw9, size=11)

txt(s9, [
    'OpenPose — PAF-based keypoints; pre-cropped hand input required for piano-scale video',
    'mido + python-rtmidi 1.5.0 — low-latency MIDI capture, delta-tick MIDI file export',
    'NumPy / Pandas — MPJPE calculations, session data structuring and aggregation',
    'Matplotlib — result visualisation (bar charts, scatter plots, regression lines)',
], In(0.5), In(4.6), W - In(1.0), In(2.6), size=12.5)


# ── slide 11 (index 10): Methodology ─────────────────────────────────────────

s11 = SL[10]
clear(s11)
ttl(s11, 'Methodology: Rapid Application Development (RAD)')

box(s11, In(0.45), In(1.0), In(7.8), In(5.1),
    'Gantt chart / project timeline\nimage placeholder',
    'Phase 1–5  ·  Weeks 1–20')

txt(s11, [
    'RAD chosen: iterative cycles suited to a solo research project',
    'Waterfall rejected: linear flow blocks mid-collection fixes',
    'Agile rejected: sprint ceremonies add overhead with no benefit',
    'Phase 1–4: hardware setup → pipeline dev → data collection → analysis',
    'Flexibility proved critical: MIDI–video sync bug resolved mid-cycle',
], In(8.55), In(1.05), In(4.3), In(5.0), size=12.5)


# ── slide 13 (index 12): Data Collection Pipeline ────────────────────────────

s13 = SL[12]
clear(s13)
ttl(s13, 'Data Collection Pipeline Overview')

pic(s13, ROOT / 'session_example.jpg',
    In(0.45), In(1.0), w=In(5.5),
    fallback_label='Session setup photograph')

box(s13, In(6.3), In(1.0), In(6.55), In(4.5),
    'Pipeline flowchart\nimage placeholder',
    'Session init → setup screen → recording → MIDI sync → output files')

txt(s13, [
    'SPACE-press sets shared rec_t0 → anchors video AND MIDI to same origin',
    'MIDI daemon thread polls every 1 ms  ·  video at ~30 FPS',
    'Output per session: .mp4  _frames.csv  _midi.jsonl  _midi.mid  _session.json',
    'No PII stored — participant identity is folder index p### only',
], In(0.5), In(5.72), W - In(1.0), In(1.55), size=12.5)


# ── slide 14 (index 13): Participant Characterisation ────────────────────────

s14 = SL[13]
clear(s14)
ttl(s14, 'Participant & Environmental Characterisation')

chars = [
    ('Fitzpatrick scale\nchart / diagram', 'Skin Tone (ITA method)'),
    ('Lux meter\nphone app screenshot', 'Ambient Illuminance'),
    ('Tape measure on\nhand photograph', 'Hand Size (cm)'),
]
pos14, cw14 = card_row(s14, chars, y_img=In(1.05), img_h=In(3.8), y_lbl=In(0), gap=In(0.4))
for (lx, _), (_, cap) in zip(pos14, chars):
    lbl(s14, cap, lx, In(4.92), cw14, size=12)

txt(s14, [
    'ITA = arctan(L*−50 / b*) × 180/π → Fitzpatrick Type I–VI  ·  sampled from palm convex hull',
    'Lux: Dim <100  ·  Indoor 100–499  ·  Bright ≥500  ·  entered manually on setup screen',
    'Hand size: palm width (cm) measured with tape before session, validated vs config bounds',
], In(0.5), In(5.48), W - In(1.0), In(1.8), size=12.5)


# ── slide 15 (index 14): Section divider — fix text only ─────────────────────

fix_mjmpe(SL[14])


# ── slide 16 (index 15): MPJPE Metric Design ─────────────────────────────────

s16 = SL[15]
clear(s16)
ttl(s16, 'MPJPE: Metric Design & Bias Reduction')

box(s16, In(0.45), In(1.0), In(5.8), In(5.1),
    'Polygon containment mask\nscreenshot placeholder',
    'Key polygons drawn over overhead camera view')

box(s16, In(6.55), In(1.0), In(6.3), In(2.35),
    'Horizontal error diagram\nh_err = |tip_x − key_cx|',
    '2-D Euclidean → x-axis only · removes finger-length bias')

txt(s16, [
    'Horizontal-only: h_err = |tip_x − key_cx| — removes depth/angle bias entirely',
    'Polygon containment: tip must lie INSIDE key boundary — 3 outcomes: matched / detection_fail / missed',
    'MPJPE accumulated only for matched events · detection_fail tracked separately',
    'Per-hand: L / R × fingers 0–4  ·  key assignment: key_cx relative to frame midpoint',
], In(6.55), In(3.55), In(6.3), In(2.6), size=12)


# ── slide 17 (index 16): Polygon & Homography ────────────────────────────────

s17 = SL[16]
clear(s17)
ttl(s17, 'Polygon Containment & Perspective Homography')

box(s17, In(0.45), In(1.0), In(6.2), In(5.1),
    'Homography calibration UI\nscreenshot placeholder',
    'Drag-corner controls · quadrilateral per key')

txt(s17, [
    'Problem: overhead camera never perfectly vertical → planar mask misaligns with real keys',
    'Solution: cv2.getPerspectiveTransform → H matrix computed once per session via drag-corner UI',
    'Warped polygon follows oblique camera angle — both models evaluated on same mask (fair comparison)',
    'Black-key perspective bias affects both models equally → comparison remains valid',
], In(6.95), In(1.05), In(5.9), In(5.0), size=12.5)


# ── slide 18 (index 17): MIDI–Video Sync ─────────────────────────────────────

s18 = SL[17]
clear(s18)
ttl(s18, 'MIDI–Video Temporal Alignment')

box(s18, In(0.45), In(1.0), In(6.2), In(5.1),
    'MIDI–video sync diagram\nimage placeholder',
    'Shared rec_t0 origin  ·  binary-search frame lookup')

txt(s18, [
    'Shared origin: rec_t0 = time.perf_counter() set at SPACE-press',
    'MidiRecorder daemon polls port.iter_pending() every 1 ms → all events timestamped to rec_t0',
    'FrameLogger records (frame_index, time_s) for every frame written → binary search at analysis time',
    'Sub-frame precision — no manual video labelling required',
    'Graceful degradation: no MIDI port found → recording continues (no crash)',
], In(6.95), In(1.05), In(5.9), In(5.0), size=12.5)


# ── slide 19 (index 18): Session Protocol ────────────────────────────────────

s19 = SL[18]
clear(s19)
ttl(s19, 'Session Protocol — Three Standard Patterns')

patterns = [
    ('Five-finger position\nC-D-E-F-G ascending\n& descending × 4', 'Pattern 1'),
    ('One-octave C major scale\nC to C with one\nfinger crossover', 'Pattern 2'),
    ('Black-key pentatonic\nC#-D#-F#-G#-A#\nascending & descending × 4', 'Pattern 3'),
]
pos19, cw19 = card_row(s19, patterns, y_img=In(1.05), img_h=In(4.2), y_lbl=In(0), gap=In(0.5))
for (lx, _), (_, cap) in zip(pos19, patterns):
    lbl(s19, cap, lx, In(5.32), cw19, size=13)

txt(s19, [
    '38 completed sessions  ·  Fitzpatrick II–V  ·  7.7–12.2 cm  ·  dim & indoor lighting',
], In(0.5), In(6.1), W - In(1.0), In(0.9), size=13, bullet=False,
   color=(180, 190, 215))


# ── slide 20 (index 19): Section divider — fix text only ─────────────────────

fix_mjmpe(SL[19])


# ── slide 21 (index 20): Dataset Overview ────────────────────────────────────

s21 = SL[20]
clear(s21)
ttl(s21, 'Dataset Overview')
pic(s21, PLOTS / '00_participant_composition.png',
    In(0.8), In(1.0), w=In(11.7),
    fallback_label='00_participant_composition.png')
txt(s21, [
    '38 sessions  ·  4 777 note-on events  ·  89 560 video frames',
    'Fitzpatrick types II–V  ·  hand sizes 7.7–12.2 cm  ·  dim & indoor lighting',
], In(0.5), In(5.85), W - In(1.0), In(1.4), size=13, bullet=False,
   color=(180, 190, 215))


# ── slide 22 (index 21): Detection Rate ──────────────────────────────────────

s22 = SL[21]
clear(s22)
ttl(s22, '8.1 Overall Detection Rate')

pic(s22, PLOTS / '03_detection_breakdown.png',
    In(0.5), In(1.0), w=In(6.1),
    fallback_label='03_detection_breakdown.png')

pic(s22, PLOTS / '14_mediapipe_per_hand_hitrate.png',
    In(6.8), In(1.0), w=In(6.0),
    fallback_label='14_mediapipe_per_hand_hitrate.png')

txt(s22, [
    'MediaPipe: 63.3% matched  ·  28.7% detection-fail  ·  ~8% missed — real-time capable (30+ FPS)',
    'OpenPose: 8.1% matched  ·  68.3% missed — requires hand-crop preprocessing; not real-time on CPU',
    'No-detection (OpenPose) loses all info; detection-fail (MediaPipe) can be snapped to nearest key',
], In(0.5), In(5.55), W - In(1.0), In(1.8), size=12.5)


# ── slide 23 (index 22): MPJPE Comparison ────────────────────────────────────

s23 = SL[22]
clear(s23)
ttl(s23, '8.2 Overall MPJPE Comparison')

pic(s23, PLOTS / '01_model_comparison_mpjpe.png',
    In(1.8), In(1.0), w=In(9.5),
    fallback_label='01_model_comparison_mpjpe.png')

txt(s23, [
    'MediaPipe: 6.07 px (n=3 017)  ·  OpenPose: 5.87 px (n=388)  ·  key half-width ≈ 10–11 px',
    'OpenPose marginally lower MPJPE but on only 8.1% of events — selection bias, not true advantage',
    'Both models below key-width threshold when detected; OpenPose unusable due to detection rate',
], In(0.5), In(5.55), W - In(1.0), In(1.8), size=12.5)


# ── slide 24 (index 23): Per-Finger MPJPE ────────────────────────────────────

s24 = SL[23]
clear(s24)
ttl(s24, '8.2.3 MediaPipe Per-Finger MPJPE & Polygon Hit Rate')

pic(s24, PLOTS / '02_per_finger_mpjpe_mediapipe.png',
    In(0.3), In(1.0), w=In(7.3),
    fallback_label='02_per_finger_mpjpe_mediapipe.png')

pic(s24, PLOTS / '14_mediapipe_per_hand_hitrate.png',
    In(7.85), In(1.0), w=In(5.0),
    fallback_label='14_mediapipe_per_hand_hitrate.png')

txt(s24, [
    'Ring & Middle: highest MPJPE — anatomically less independent digits',
    'Thumb & Index: lowest MPJPE — most motor-independent',
    'Right hand 95.8% hit rate vs Left 53.2% — right-hand dominance artifact, not model failure',
    'Despite asymmetric hit rate: Left 6.13 px vs Right 6.02 px — near-identical MPJPE',
], In(0.5), In(5.55), W - In(1.0), In(1.8), size=12.5)


# ── slide 25 (index 24): Skin Tone ───────────────────────────────────────────

s25 = SL[24]
clear(s25)
ttl(s25, '8.3 Effect of Skin Tone (Fitzpatrick Type) on MPJPE')

pic(s25, PLOTS / '05_mpjpe_by_fitzpatrick.png',
    In(0.3), In(1.0), w=In(6.3),
    fallback_label='05_mpjpe_by_fitzpatrick.png')

pic(s25, PLOTS / '09_detection_fail_by_fitzpatrick.png',
    In(6.85), In(1.0), w=In(6.0),
    fallback_label='09_detection_fail_by_fitzpatrick.png')

txt(s25, [
    'MediaPipe MPJPE range: Type II 5.51 px → Type V 6.29 px  ·  total range 0.88 px (within SD)',
    'Trend not monotonic → no systematic skin-tone bias across tested range',
    'Detection-fail: Type III 32.8% (highest) vs Type II 16.8% — no consistent direction',
    'Types I & VI not represented — unable to test extremes of Fitzpatrick scale',
], In(0.5), In(5.55), W - In(1.0), In(1.8), size=12.5)


# ── slide 26 (index 25): Lighting ────────────────────────────────────────────

s26 = SL[25]
clear(s26)
ttl(s26, '8.4 Effect of Ambient Illuminance on MPJPE')

pic(s26, PLOTS / '04_mpjpe_by_lux.png',
    In(0.3), In(1.0), w=In(6.1),
    fallback_label='04_mpjpe_by_lux.png')

pic(s26, PLOTS / '16_mpjpe_vs_lux_scatter.png',
    In(6.7), In(1.0), w=In(6.3),
    fallback_label='16_mpjpe_vs_lux_scatter.png')

txt(s26, [
    'MediaPipe: Dim 6.03 px (SD=1.16) vs Indoor 6.17 px (SD=0.76) — 0.14 px gap, within SD',
    'Regression nearly flat across full lux range — lighting does not affect MediaPipe accuracy',
    'Wide spread within dim band → individual skill is the dominant variance source, not lighting',
    'No bright (>500 lux) sessions recorded — gap for future research',
], In(0.5), In(5.55), W - In(1.0), In(1.8), size=12.5)


# ── slide 27 (index 26): Hand Size ───────────────────────────────────────────

s27 = SL[26]
clear(s27)
ttl(s27, '8.5 Effect of Hand Size on MPJPE')

pic(s27, PLOTS / '06_mpjpe_vs_handsize.png',
    In(1.8), In(1.0), w=In(9.5),
    fallback_label='06_mpjpe_vs_handsize.png')

txt(s27, [
    'MediaPipe: slope 0.162 px/cm  ·  R² = 0.030 — total variation ~0.7 px across full range',
    'OpenPose: slope 0.336 px/cm  ·  R² = 0.059 — also negligible (small n caveat applies)',
    'Hand size explains <6% of MPJPE variance — no per-user hand-size calibration needed',
], In(0.5), In(5.55), W - In(1.0), In(1.8), size=12.5)


# ── slide 29 (index 28): Novel Contributions — fix MPJPE text ────────────────

fix_mjmpe(SL[28])


# ── slide 30 (index 29): Conclusion — fix MPJPE text ─────────────────────────

fix_mjmpe(SL[29])


# ── slide 31 (index 30): Future Work ─────────────────────────────────────────

s31 = SL[30]
clear(s31)
ttl(s31, 'Future Work')

box(s31, In(0.45), In(1.05), In(4.5), In(5.6),
    'Future research directions\nimage placeholder',
    'e.g. Fitzpatrick I/VI study  ·  Kalman filtering  ·  3-D mask')

txt(s31, [
    'Recruit Fitzpatrick Type I & VI participants — confirm equitable performance at extremes',
    'Include participants under 18 (guardian consent) — hand sizes < 7.7 cm',
    'Extend lighting range: outdoor bright >500 lux, very dark <5 lux',
    'Add OpenPose hand-crop preprocessing for fair head-to-head comparison',
    'Kalman / optical-flow interpolation to reduce 28.7% MediaPipe detection-fail rate',
    'Recruit experienced pianists for complex repertoire — occlusion stress-test',
    'Replace 2-D polygon mask with per-key depth-corrected mask for upper-register accuracy',
], In(5.35), In(1.05), In(7.5), In(5.6), size=12.5)


# ── save ──────────────────────────────────────────────────────────────────────

prs.save(PPTX)
print('Saved:', PPTX)
